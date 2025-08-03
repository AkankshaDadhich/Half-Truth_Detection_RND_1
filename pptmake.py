from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Updated classification report data  
data = [
    ["Model", "Accuracy", "Macro F1-Score", "Macro Recall", "Macro Precision"],
    ["SVM", "0.78", "0.78", "0.78", "0.79"],
    ["Random Forest", "0.68", "0.67", "0.68", "0.68"],
    ["Decision Tree", "0.70", "0.70", "0.70", "0.78"],
    ["Gradient Boosting Classifier", "0.70", "0.70", "0.70", "0.70"]
]

# Create a PowerPoint presentation and add a slide
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

# Add title to the slide
title = slide.shapes.title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
title_frame = title.text_frame
title_frame.text = "Classification Report"
title_frame.paragraphs[0].font.size = Pt(24)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Define table dimensions
rows, cols = len(data), len(data[0])
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3))
table = table_shape.table

# Populate table data
for row_idx, row_data in enumerate(data):
    for col_idx, cell_data in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_data
        # Set font size and alignment
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)  # Adjust font size
            paragraph.font.bold = True if row_idx == 0 else False  # Bold header row
            paragraph.alignment = PP_ALIGN.CENTER
            if row_idx == 0:  # Apply styling to header row
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 102, 204)

# Save the PowerPoint file
prs.save("classification_report_ppt_table.pptx")
print("PowerPoint file saved as 'classification_report_ppt_table.pptx'")
